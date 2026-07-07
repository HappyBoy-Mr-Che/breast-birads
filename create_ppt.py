"""
Generate PPT v4 — robust flowcharts, spread visualization images, analysis text.
Uses only RECTANGLE + text boxes for maximum PowerPoint compatibility.
"""
import json, os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════ Paths & Data ═══════════════════
PROJECT_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = PROJECT_DIR / "new_tools" / "outputs"
LOG_PATH = OUTPUT_DIR / "logs" / "training_log.json"
IMG_DIR = OUTPUT_DIR / "logs"
PPT_PATH = Path(os.path.expanduser("~/Desktop")) / "乳腺超声BI-RADS多任务算法汇报_v5.pptx"

HAS_DATA = False; best_epoch = None; best_f1 = 0.0; total_epochs = 0
try:
    with open(LOG_PATH) as f:
        train_log = json.load(f)
    total_epochs = len(train_log)
    for d in train_log:
        f1k = ["boundary_f1","calcification_f1","shape_f1","direction_f1"]
        mf1 = sum(d["val"].get(k,0) for k in f1k)/4.0
        if mf1 > best_f1: best_f1 = mf1; best_epoch = d
    HAS_DATA = True
except Exception as e: print(f"Warning: {e}")

# ═══════════════════ Colors ═══════════════════
PRI   = RGBColor(0x0D,0x2B,0x4E)
SEC   = RGBColor(0x1A,0x56,0x8C)
ACC   = RGBColor(0x2E,0x86,0xC1)
HIL   = RGBColor(0xE7,0x6F,0x51)
GRN   = RGBColor(0x27,0xAE,0x60)
WHT   = RGBColor(0xFF,0xFF,0xFF)
BG    = RGBColor(0xF8,0xF9,0xFA)
G100  = RGBColor(0xF1,0xF3,0xF5)
G200  = RGBColor(0xDE,0xE2,0xE6)
G300  = RGBColor(0xCE,0xD4,0xDA)
G400  = RGBColor(0xAD,0xB5,0xBD)
G600  = RGBColor(0x6C,0x75,0x7D)
G800  = RGBColor(0x34,0x3A,0x40)
TXT   = RGBColor(0x21,0x25,0x29)

SW = Inches(13.333); SH = Inches(7.5)
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

# ═══════════════════ Helpers ═══════════════════
def BS(): return prs.slides.add_slide(prs.slide_layouts[6])
def BG(s,c=BG): s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def R(s,l,t,w,h,fill=None,border=None,bw=Pt(1)):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else: sp.fill.background()
    if border: sp.line.color.rgb = border; sp.line.width = bw
    else: sp.line.fill.background()
    return sp
def T(s,l,t,w,h,txt="",sz=14,c=TXT,b=False,a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l,t,w,h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = "Microsoft YaHei"; p.alignment = a
    return tb
def MT(s,l,t,w,h,lines):
    tb = s.shapes.add_textbox(l,t,w,h); tf = tb.text_frame; tf.word_wrap = True
    for i,item in enumerate(lines):
        txt=item[0]; sz=item[1]if len(item)>1 else 12; clr=item[2]if len(item)>2 else TXT
        bd=item[3]if len(item)>3 else False; al=item[4]if len(item)>4 else PP_ALIGN.LEFT
        sp=item[5]if len(item)>5 else Pt(4)
        p=tf.paragraphs[0]if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=clr; p.font.bold=bd
        p.font.name="Microsoft YaHei"; p.alignment=al; p.space_after=sp
    return tb
def TB(s,title,subtitle=None):
    R(s,Inches(0),Inches(0),SW,Inches(0.05),fill=HIL)
    R(s,Inches(0),Inches(0.05),SW,Inches(1.05),fill=PRI)
    T(s,Inches(0.7),Inches(0.1),Inches(12),Inches(0.6),title,sz=30,c=WHT,b=True)
    if subtitle: T(s,Inches(0.7),Inches(0.65),Inches(12),Inches(0.35),subtitle,sz=14,c=G400)
    R(s,Inches(0),Inches(1.1),SW,Inches(0.03),fill=G200)
def PN(s,n): T(s,SW-Inches(1.0),SH-Inches(0.45),Inches(0.8),Inches(0.35),f"{n}/{TOTAL}",sz=9,c=G400,a=PP_ALIGN.RIGHT)
def ST(s,l,t,txt,sz=20):
    R(s,l,t+Inches(0.05),Inches(0.06),Inches(0.4),fill=ACC)
    T(s,l+Inches(0.2),t,Inches(8),Inches(0.45),txt,sz=sz,c=PRI,b=True)
def AR(s,l,t,w,h,c=G400):
    sp=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background()
    return sp
def DA(s,l,t,w,h,c=G400):
    sp=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background()
    return sp
def card(s,l,t,w,h,fill=WHT,border=G200): return R(s,l,t,w,h,fill=fill,border=border)
def tbl(s,data,l,t,w,h,cw,hd=PRI):
    rows=len(data); cols=len(data[0])
    ts=s.shapes.add_table(rows,cols,l,t,w,h); tb=ts.table
    for ci,cww in enumerate(cw): tb.columns[ci].width=cww
    for r in range(rows):
        for c in range(cols):
            cell=tb.cell(r,c); cell.text=str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.name="Microsoft YaHei"; p.font.size=Pt(13); p.alignment=PP_ALIGN.CENTER
                if r==0: p.font.bold=True; p.font.color.rgb=WHT
                else: p.font.color.rgb=G800
            if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=hd
            elif r%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=G100
    return ts

# ── Flowchart helper: draws a row of labeled boxes with arrows ──
def flow_row(s, x0, y, items, bw=Inches(2.3), bh=Inches(1.6), gap=Inches(0.15)):
    """items: list of (title, subtitle, color). Draws a horizontal flowchart row."""
    for i,(t,st,clr) in enumerate(items):
        ix=x0+i*(bw+gap)
        R(s,ix,y,bw,bh,fill=clr)
        T(s,ix+Inches(0.08),y+Inches(0.1),bw-Inches(0.16),Inches(0.55),t,sz=11,c=WHT,b=True,a=PP_ALIGN.CENTER)
        T(s,ix+Inches(0.08),y+Inches(0.7),bw-Inches(0.16),Inches(0.8),st,sz=9.5,c=RGBColor(0xEE,0xEE,0xFF),a=PP_ALIGN.CENTER)
        if i<len(items)-1: AR(s,ix+bw+Inches(0.02),y+Inches(0.65),Inches(0.12),Inches(0.18))

def img_with_analysis(s, img_fn, ix, iy, iw, ih, title, analysis_lines):
    """Add an image card with analysis text to the right."""
    card(s,ix,iy,iw,ih)
    T(s,ix+Inches(0.12),iy+Inches(0.06),iw-Inches(0.24),Inches(0.28),title,sz=14,c=PRI,b=True)
    img_path = IMG_DIR / img_fn
    img_ok = False
    if img_path.exists():
        try:
            s.shapes.add_picture(str(img_path),ix+Inches(0.1),iy+Inches(0.4),iw-Inches(0.2),ih-Inches(0.5))
            img_ok = True
        except: pass
    if analysis_lines:
        # Analysis text below or to the right
        pass
    return img_ok

TOTAL = 18

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s,PRI)
R(s,Inches(0),Inches(0),SW,Inches(0.07),fill=HIL)
tri=s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,Inches(9.5),Inches(0),Inches(3.833),Inches(5.3))
tri.fill.solid(); tri.fill.fore_color.rgb=RGBColor(0x15,0x3A,0x60); tri.line.fill.background()
T(s,Inches(1.2),Inches(2.0),Inches(9),Inches(1.0),"乳腺超声 BI-RADS 分类与特征检测",sz=42,c=WHT,b=True)
T(s,Inches(1.2),Inches(3.0),Inches(9),Inches(0.6),"多任务深度学习算法 — 技术汇报",sz=24,c=RGBColor(0x9B,0xB7,0xD4))
R(s,Inches(1.2),Inches(3.8),Inches(1.0),Inches(0.05),fill=HIL)
T(s,Inches(1.2),Inches(4.1),Inches(9),Inches(0.5),"Multi-Task Deep Learning for Breast Ultrasound BI-RADS Classification & Feature Detection",sz=14,c=RGBColor(0x7B,0x97,0xB4))
T(s,Inches(1.2),Inches(5.0),Inches(9),Inches(0.4),"EfficientNet-B3 + FPN Neck + 5-Head Multi-Task Architecture",sz=15,c=G600)
T(s,Inches(1.2),Inches(5.8),Inches(9),Inches(0.4),"2025",sz=18,c=G600)
PN(s,1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview + Architecture
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"项目概述与技术总览","Project Overview & Multi-Task Architecture")

# Left: Problem
ST(s,Inches(0.6),Inches(1.4),"临床问题定义")
prob=[
("■ 乳腺癌是全球女性发病率最高的恶性肿瘤，乳腺超声是首选影像筛查手段",14,TXT,False,PP_ALIGN.LEFT,Pt(4)),
("■ BI-RADS 是国际标准化的乳腺影像评估分级系统：2类->3类->4A类->4B类->4C类->5类",14,TXT,False,PP_ALIGN.LEFT,Pt(6)),
("■ 4项关键超声特征对良恶性判断至关重要：边界光滑度·钙化存在性·形态规则性·生长方向",14,SEC,True,PP_ALIGN.LEFT,Pt(6)),
("■ 传统人工判读主观性强、耗时长，亟需AI辅助诊断",14,TXT,False),
]
MT(s,Inches(0.6),Inches(2.0),Inches(5.8),Inches(2.5),prob)

# Right: Objectives
ST(s,Inches(6.8),Inches(1.4),"研究目标 — 三项任务统一框架")
obj=[
("构建一个多任务深度学习模型，同时完成：",14,TXT,False),
("",6,TXT,False),
("1.  BI-RADS 6类分级诊断",16,ACC,True),
("    判定病灶的BI-RADS等级 (2类/3类/4A/4B/4C/5类)",13,G800,False),
("",6,TXT,False),
("2.  4项超声特征检测 + 区域定位",16,ACC,True),
("    边界(光滑/不光滑) | 钙化(无/有) | 形状(规则/不规则) | 方向(平行/不平行)",13,G800,False),
("    每个特征同时输出：二分类标签 + 边界框 [xc, yc, w, h]",13,G800,False),
]
MT(s,Inches(6.8),Inches(2.0),Inches(6.0),Inches(2.5),obj)

# Architecture diagram — compact, within bounds
arch_y=Inches(4.5)
# Input
R(s,Inches(5.2),arch_y,Inches(3.0),Inches(0.42),fill=PRI)
T(s,Inches(5.25),arch_y+Inches(0.04),Inches(2.9),Inches(0.34),"输入图像 (B, 3, 640x320)",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)
DA(s,Inches(6.5),arch_y+Inches(0.42),Inches(0.22),Inches(0.15))
# Backbone
R(s,Inches(3.5),arch_y+Inches(0.7),Inches(6.4),Inches(0.42),fill=SEC)
T(s,Inches(3.55),arch_y+Inches(0.72),Inches(6.3),Inches(0.38),"EfficientNet-B3 骨干网络  (ImageNet预训练)  |  1536 ch  |  ~12.2M参数",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)
DA(s,Inches(6.5),arch_y+Inches(1.12),Inches(0.22),Inches(0.15))
# FPN Neck
R(s,Inches(4.5),arch_y+Inches(1.4),Inches(4.4),Inches(0.38),fill=ACC)
T(s,Inches(4.55),arch_y+Inches(1.42),Inches(4.3),Inches(0.34),"SimpleNeck (FPN)  多尺度融合 -> 256 ch",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)

# 5 heads in a row
heads_y=arch_y+Inches(2.0)
heads=[("BI-RADS\n分类头","6类",PRI),("边界\n检测头","2类+BBox",SEC),("钙化\n检测头","2类+BBox",SEC),("形状\n检测头","2类+BBox",SEC),("方向\n检测头","2类+BBox",SEC)]
hw=Inches(2.4);hg=Inches(0.12);total_w=5*hw+4*hg;hx0=(SW-total_w)/2
for i,(nm,out,clr) in enumerate(heads):
    hx=hx0+i*(hw+hg)
    R(s,hx,heads_y,hw,Inches(0.6),fill=clr)
    T(s,hx+Inches(0.05),heads_y+Inches(0.04),hw-Inches(0.1),Inches(0.3),nm,sz=10,c=WHT,b=True,a=PP_ALIGN.CENTER)
    T(s,hx+Inches(0.05),heads_y+Inches(0.34),hw-Inches(0.1),Inches(0.24),out,sz=9,c=RGBColor(0xDD,0xEE,0xFF),a=PP_ALIGN.CENTER)
# Output row
out_y=heads_y+Inches(0.7)
outputs=["2类/3类/4A类\n4B类/4C类/5类","光滑/不光滑\n[xc,yc,w,h]","无钙化/有钙化\n[xc,yc,w,h]","规则/不规则\n[xc,yc,w,h]","平行/不平行\n[xc,yc,w,h]"]
for i,(nm,out,clr) in enumerate(heads):
    hx=hx0+i*(hw+hg)
    R(s,hx+Inches(0.15),out_y,hw-Inches(0.3),Inches(0.65),fill=WHT,border=clr,bw=Pt(1.5))
    T(s,hx+Inches(0.2),out_y+Inches(0.05),hw-Inches(0.4),Inches(0.55),outputs[i],sz=9,c=clr,b=True,a=PP_ALIGN.CENTER)

PN(s,2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Data Processing Pipeline (robust flowchart)
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"数据处理全流程","Data Processing Pipeline — from Raw Data to Model Input")

# Flow row 1 — 5 steps
flow_row(s, Inches(0.4), Inches(1.55), [
    ("数据采集与组织","两个数据源:\nfuture/ + classfy/\n6个BI-RADS类别\n4种特征标签\nYOLO格式标注",PRI),
    ("数据预处理\npreprocess.py","统计分析\n完整性校验\n损坏图像检测\n85/15划分\nJSON索引文件",SEC),
    ("数据加载\nDataset类","融合双数据源\nYOLO标签解析\n特征标签映射\n缺失标签默认值\n类别映射构建",SEC),
    ("数据增强\nTransforms类","随机仿射(+-20deg)\n水平翻转(50%)\n亮度对比度调整\n弹性变形(30%)\nMixUp/CutMix",ACC),
    ("归一化&批处理\nCollate函数","Resize->640x320\nImageNet标准化\nmean/std归一化\nTensor组装\n(B,3,640,320)",HIL),
])

# Two data source panels
ST(s,Inches(0.6),Inches(3.6),"数据源详细结构")
# Left panel
card(s,Inches(0.6),Inches(4.2),Inches(5.8),Inches(2.9))
R(s,Inches(0.6),Inches(4.2),Inches(5.8),Inches(0.5),fill=SEC)
T(s,Inches(0.75),Inches(4.25),Inches(5.5),Inches(0.4),"特征检测数据 (future/train/) — 4种超声特征",sz=14,c=WHT,b=True)
src1=[
("images/               统一的超声图像存储目录 (.jpg/.png)",13,G800,False),
("boundary_labels/       边界光滑度: smooth / not_smooth",13,G800,False),
("calcification_labels/  钙化存在性: no_calcification / calcification",13,G800,False),
("shape_labels/          形态规则性: regular / irregular",13,G800,False),
("direction_labels/      生长方向: parallel / not_parallel",13,G800,False),
("",8,TXT,False),
("标签格式: YOLO格式 — class xc yc w h (归一化到[0,1])",13,ACC,True),
("缺失标签的样本: 默认填充 [0, 0.0, 0.0, 0.0, 0.0] (负样本+零bbox)",12,G600,False),
]
MT(s,Inches(0.75),Inches(4.8),Inches(5.5),Inches(2.2),src1)

# Right panel
card(s,Inches(6.9),Inches(4.2),Inches(5.8),Inches(2.9))
R(s,Inches(6.9),Inches(4.2),Inches(5.8),Inches(0.5),fill=ACC)
T(s,Inches(7.05),Inches(4.25),Inches(5.5),Inches(0.4),"BI-RADS分类数据 (classfy/train/) — 6个等级",sz=14,c=WHT,b=True)
src2=[
("2类/images/     BI-RADS 2类 — 良性发现",13,G800,False),
("3类/images/     BI-RADS 3类 — 可能良性（建议短期随访）",13,G800,False),
("4A类/images/    BI-RADS 4A类 — 低度可疑恶性",13,G800,False),
("4B类/images/    BI-RADS 4B类 — 中度可疑恶性",13,G800,False),
("4C类/images/    BI-RADS 4C类 — 高度可疑恶性",13,G800,False),
("5类/images/     BI-RADS 5类 — 高度提示恶性",13,G800,False),
("",8,TXT,False),
("数据组织: 按类别文件夹分存, 每个类别下的 images/ 存放对应超声图像",13,ACC,True),
("映射方式: 文件夹名->类别索引, 建立 stem_to_birads 全局映射表",12,G600,False),
]
MT(s,Inches(7.05),Inches(4.8),Inches(5.5),Inches(2.2),src2)

PN(s,3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Data Augmentation (robust flowchart)
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"数据增强策略详解","Data Augmentation Pipeline — Transforms Class Implementation")

# Flow row — 7 steps
flow_row(s, Inches(0.3), Inches(1.4), [
    ("原始图像\nH x W x 3","乳腺超声\n原始采集图像",PRI),
    ("随机仿射变换","旋转 +-20deg\n缩放 0.85-1.15\n平移 +-5%\nbbox同步变换",SEC),
    ("水平翻转","概率 p=0.5\n图像水平镜像\nbbox xc -> 1-xc",SEC),
    ("亮度/对比度","alpha=1+-0.3\nbeta=+-30\n30%概率附加\ngamma校正(LUT)",SEC),
    ("弹性变形","alpha=40 sigma=4\n概率 p=0.3\n高斯滤波平滑\ncv2.remap重映射",ACC),
    ("Resize","640 x 320\nBILINEAR插值\n统一输入尺寸",ACC),
    ("归一化","ImageNet统计量\nmean/std\n除以255\n-> Tensor(3,H,W)",HIL),
], bw=Inches(1.65), bh=Inches(1.45), gap=Inches(0.1))

# Three method detail cards
ST(s,Inches(0.6),Inches(3.3),"增强方法技术细节")
methods=[
("随机仿射变换\n(RandomAffine)","核心: cv2.getRotationMatrix2D 生成复合变换矩阵\n\n参数配置:\n  旋转: uniform(-20deg, +20deg)\n  缩放: uniform(0.85, 1.15)\n  平移: uniform(-5%, +5%) x 图像尺寸\n\n关键实现: bbox坐标通过矩阵乘法同步变换\n  保证增强后标签与图像严格对应\n\n边界模式: BORDER_REFLECT (镜像填充)\n  避免边缘产生人工伪影",SEC),
("弹性变形\n(ElasticTransform)","目的: 模拟超声探头压力导致的组织形变\n\n实现步骤:\n  1. 生成 (H,W) 随机位移场 (均匀分布)\n  2. 高斯滤波平滑 (sigma=4)\n  3. cv2.remap 执行像素级重映射\n\n参数:\n  alpha=40  控制变形幅度\n  sigma=4   控制平滑程度\n  p=0.3     仅30%概率执行\n\n注意: 仅作用于图像本身\n  不影响分类和检测标签",ACC),
("MixUp / CutMix\n批次级混合增强","训练时 50% 步骤启用, 作为强正则化\n\nMixUp (线性混合):\n  img = lambda * img_A + (1-lambda) * img_B\n  lambda ~ Beta(0.2, 0.2), 保证 lambda >= 0.5\n  标签和损失函数同步按 lambda 加权\n\nCutMix (区域替换):\n  从图像B随机切出矩形区域\n  粘贴到图像A对应位置\n  面积比例即混合系数\n\n效果: 显著提升泛化, 减少过拟合",HIL),
]
for i,(title,desc,clr) in enumerate(methods):
    mx=Inches(0.5+i*4.15)
    card(s,mx,Inches(3.8),Inches(3.95),Inches(3.4))
    R(s,mx,Inches(3.8),Inches(3.95),Inches(0.5),fill=clr)
    T(s,mx+Inches(0.1),Inches(3.85),Inches(3.75),Inches(0.4),title,sz=13,c=WHT,b=True)
    MT(s,mx+Inches(0.1),Inches(4.4),Inches(3.75),Inches(2.7),
       [(line,10.5,G800 if line else TXT,False,PP_ALIGN.LEFT,Pt(2)if line else Pt(4))for line in desc.split("\n")])

PN(s,4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Backbone (redesigned: table-based architecture + large flow)
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"骨干网络详解 — EfficientNet-B3","Backbone Architecture — EfficientNet-B3 Deep Dive")

# ── Left panel: Architecture as a clear TABLE (guaranteed to render) ──
ST(s,Inches(0.5),Inches(1.4),"EfficientNet-B3 层级结构")

arch_table = [
    ["Stage", "模块类型", "配置", "输出通道", "分辨率"],
    ["0", "Conv3x3, stride=2", "BN + SiLU", "40", "320x160"],
    ["1", "MBConv1, k3x3", "SE, x1", "24", "160x80"],
    ["2", "MBConv6, k3x3", "SE, x2", "32", "80x40"],
    ["3", "MBConv6, k5x5", "SE, x2", "48", "40x20"],
    ["4", "MBConv6, k3x3", "SE, x3", "96", "20x10"],
    ["5", "MBConv6, k5x5", "SE, x4", "136", "20x10"],
    ["6", "MBConv6, k5x5", "SE, x5", "232", "20x10"],
    ["7", "Conv1x1 + Pool + FC", "—", "1536", "1x1"],
]
tbl(s, arch_table, Inches(0.5), Inches(2.0), Inches(7.0), Inches(3.8),
    [Inches(0.6), Inches(1.6), Inches(1.1), Inches(0.9), Inches(0.9)])
# Note below table
T(s, Inches(0.5), Inches(5.9), Inches(7.0), Inches(0.4),
   "MBConv: Mobile Inverted Bottleneck Convolution (内置 SE 注意力) | 总层数 ~264 层 | 参数量 ~12.2M | ImageNet Top-1: 84.1%",
   sz=12, c=PRI, b=True)

# Architecture flow diagram — simple horizontal blocks, uniform height
ST(s, Inches(0.5), Inches(6.3), "数据流向")
flow_blocks = [
    ("Input\n640x320x3", PRI),
    ("Stage 0\nConv", SEC),
    ("Stage 1-3\nMBConv", SEC),
    ("Stage 4-6\nMBConv", ACC),
    ("Stage 7\nFC", PRI),
    ("Output\n1536 ch\n20x10", HIL),
]
fb_w = Inches(1.75); fb_h = Inches(0.8); fb_gap = Inches(0.15)
fb_x0 = Inches(0.5)
for i, (label, clr) in enumerate(flow_blocks):
    fx = fb_x0 + i*(fb_w+fb_gap)
    R(s, fx, Inches(6.45), fb_w, fb_h, fill=clr)
    T(s, fx+Inches(0.05), Inches(6.48), fb_w-Inches(0.1), fb_h-Inches(0.1),
       label, sz=10, c=WHT, b=True, a=PP_ALIGN.CENTER)
    if i < len(flow_blocks)-1:
        AR(s, fx+fb_w+Inches(0.02), Inches(6.65), Inches(0.12), Inches(0.18))

# ── Right panel: Comparison + Selection ──
# Comparison table
ST(s, Inches(7.8), Inches(1.4),"候选骨干网络对比")
bbd=[
["骨干网络","输出ch","参数量","FPN","Top-1","速度"],
["EfficientNet-B3","1536","12.2M","单尺度","84.1%","快"],
["ResNet-50","2048","25.6M","多尺度","80.9%","中"],
["ResNet-34","512","21.8M","多尺度","73.3%","快"],
["ConvNeXt-Tiny","768","28.6M","单尺度","82.1%","中"],
["MobileNetV3-L","960","5.5M","单尺度","75.2%","极快"],
]
tbl(s,bbd,Inches(7.8),Inches(2.0),Inches(5.1),Inches(2.6),
    [Inches(1.45),Inches(0.7),Inches(0.7),Inches(0.65),Inches(0.65),Inches(0.6)])

# Selection rationale
ST(s,Inches(7.8),Inches(4.9),"选择 EfficientNet-B3 的五大理由")
why=[
("1. 精度-效率最优: 仅12M参数即达84.1% Top-1 Acc",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("2. SE注意力机制: MBConv内置Squeeze-and-Excitation，自适应通道加权",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("3. 感受野匹配: B3规模(输入640x320)恰好匹配乳腺病灶的典型空间尺度",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("4. 迁移学习效果好: ImageNet预训练权重对超声纹理模式有良好泛化能力",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("5. 训练效率高: 相比ResNet-50(25.6M)，参数减半，训练吞吐量更高",13,G800,False),
]
MT(s,Inches(7.8),Inches(5.35),Inches(5.1),Inches(1.8),why)

PN(s,5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — FPN Neck + Classification Head
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"FPN颈部与分类头 — 算法结构详解","SimpleNeck FPN & ClassificationHead Architecture")

# Left: FPN
ST(s,Inches(0.6),Inches(1.4),"SimpleNeck — 轻量级多尺度特征融合颈部")
card(s,Inches(0.6),Inches(1.95),Inches(6.0),Inches(5.2))

fx=Inches(0.8); fy=Inches(2.1)
# Input -> Compress
R(s,fx+Inches(1.5),fy,Inches(2.5),Inches(0.45),fill=PRI)
T(s,fx+Inches(1.55),fy+Inches(0.05),Inches(2.4),Inches(0.35),"骨干输出: 1536 ch, 20x10",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)
DA(s,fx+Inches(2.55),fy+Inches(0.45),Inches(0.22),Inches(0.16))
R(s,fx+Inches(1.0),fy+Inches(0.73),Inches(3.5),Inches(0.4),fill=SEC)
T(s,fx+Inches(1.05),fy+Inches(0.75),Inches(3.4),Inches(0.35),"Conv2d(1536->256, k=1) + BN + ReLU",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)

# Three branches
br_y=fy+Inches(1.35)
R(s,fx,br_y,Inches(1.8),Inches(0.7),fill=ACC)
T(s,fx+Inches(0.05),br_y+Inches(0.05),Inches(1.7),Inches(0.6),"0.5x 尺度\nConv3x3 stride=2\n感受野: 粗粒度(全局)",sz=10,c=WHT,b=True,a=PP_ALIGN.CENTER)
R(s,fx+Inches(2.05),br_y,Inches(1.8),Inches(0.7),fill=SEC)
T(s,fx+Inches(2.1),br_y+Inches(0.05),Inches(1.7),Inches(0.6),"1x (原始) 尺度\nConv3x3 pad=1\n感受野: 中粒度",sz=10,c=WHT,b=True,a=PP_ALIGN.CENTER)
R(s,fx+Inches(4.1),br_y,Inches(1.8),Inches(0.7),fill=ACC)
T(s,fx+Inches(4.15),br_y+Inches(0.05),Inches(1.7),Inches(0.6),"2x 尺度\nUpsample Bilinear\n+ Conv3x3 感受野: 细粒度",sz=10,c=WHT,b=True,a=PP_ALIGN.CENTER)
T(s,fx,br_y+Inches(0.75),Inches(6.0),Inches(0.22),"均双线性插值回原始分辨率 20x10",sz=10,c=G600,a=PP_ALIGN.CENTER)

# Fusion
R(s,fx+Inches(1.0),br_y+Inches(1.1),Inches(3.8),Inches(0.4),fill=HIL)
T(s,fx+Inches(1.05),br_y+Inches(1.12),Inches(3.7),Inches(0.35),"Concat(256x3=768) -> Conv1x1 -> 256 ch 融合输出",sz=12,c=WHT,b=True,a=PP_ALIGN.CENTER)

fpn_adv=[
("设计思想:",13,PRI,True,PP_ALIGN.LEFT,Pt(2)),
("  3个不同感受野并行处理同一特征图，1x1卷积自动学习最优融合权重",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("  极轻量 (~590K参数) | 不改变空间分辨率 | 兼容所有单尺度骨干网络",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("  检测头获得粗/中/细三种粒度的空间上下文，提升定位精度",12,G800,False),
]
MT(s,fx,br_y+Inches(1.7),Inches(5.8),Inches(1.5),fpn_adv)

# Right: Classification Head
ST(s,Inches(7.0),Inches(1.4),"ClassificationHead — BI-RADS 6类分类头")
card(s,Inches(7.0),Inches(1.95),Inches(5.8),Inches(5.2))

cls_layers=[
("AdaptiveAvgPool2d(1x1) -> Flatten","空间压缩 -> (B, 1536)",PRI),
("Dropout(0.50)   高比例正则化","防止骨干特征过拟合",SEC),
("Linear(1536->512) + BatchNorm1d + ReLU","FC1: 1536 -> 512",SEC),
("Dropout(0.30)   中等比例正则化","逐层递减Dropout策略",SEC),
("Linear(512->256) + BatchNorm1d + ReLU","FC2: 512 -> 256",ACC),
("Dropout(0.15)   低比例正则化","渐进式降低正则化强度",ACC),
("Linear(256->6)   输出层","6类 BI-RADS logits",HIL),
]
cly=Inches(2.1)
for i,(layer,note,clr) in enumerate(cls_layers):
    R(s,Inches(7.15),cly+i*Inches(0.7),Inches(5.5),Inches(0.6),fill=clr)
    T(s,Inches(7.25),cly+i*Inches(0.7)+Inches(0.02),Inches(3.5),Inches(0.28),layer,sz=11,c=WHT,b=True)
    T(s,Inches(7.25),cly+i*Inches(0.7)+Inches(0.3),Inches(3.5),Inches(0.25),note,sz=10,c=RGBColor(0xDD,0xEE,0xFF))

ds_y=cly+7*Inches(0.7)+Inches(0.2)
design=[
("设计要点:",14,PRI,True,PP_ALIGN.LEFT,Pt(3)),
("  - 逐层递减 Dropout (0.5 -> 0.3 -> 0.15):",12,G800,False),
("    浅层强正则化防过拟合，语义层保留更多信息",12,G800,False),
("  - BatchNorm1d: 加速收敛，允许更高学习率",12,G800,False),
("  - 3层MLP: 充分非线性 + 避免过深导致过拟合",12,G800,False),
]
MT(s,Inches(7.15),ds_y,Inches(5.5),Inches(1.5),design)

PN(s,6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Feature Detection Head
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"特征检测头 — 算法结构详解","FeatureDetectionHead — Dual-Branch Architecture (x4)")

ST(s,Inches(0.6),Inches(1.4),"FeatureDetectionHead 内部结构 (4个特征各自独立，结构完全相同)")

# Conv stage
cv_y=Inches(2.1)
cv_steps=[
("Conv2d(256->256,k3,p1)\nBN + ReLU + Drop2d(0.2)","保持空间分辨率\nDrop2d抑制空间过拟合",PRI),
("Conv2d(256->256,k3,p1)\nBN + ReLU + Drop2d(0.2)","进一步精炼特征\n增强判别能力",SEC),
("Conv2d(256->128,k3,p1)\nBN + ReLU (无Drop2d)","压缩到128通道\n为FC层做准备",ACC),
]
for i,(desc,note,clr) in enumerate(cv_steps):
    cx=Inches(0.6+i*4.2)
    R(s,cx,cv_y,Inches(3.8),Inches(1.0),fill=clr)
    T(s,cx+Inches(0.1),cv_y+Inches(0.05),Inches(3.6),Inches(0.55),desc,sz=11,c=WHT,b=True,a=PP_ALIGN.CENTER)
    T(s,cx+Inches(0.1),cv_y+Inches(0.62),Inches(3.6),Inches(0.3),note,sz=9.5,c=RGBColor(0xDD,0xEE,0xFF),a=PP_ALIGN.CENTER)
    if i<2: AR(s,cx+Inches(3.8)+Inches(0.05),cv_y+Inches(0.35),Inches(0.16),Inches(0.18))

# Pool
pool_y=cv_y+Inches(1.3)
R(s,Inches(2.0),pool_y,Inches(9.3),Inches(0.55),fill=HIL)
T(s,Inches(2.05),pool_y+Inches(0.08),Inches(9.2),Inches(0.4),"AdaptiveAvgPool2d(1x1) -> Flatten -> Dropout(0.5) -> 128维特征向量",sz=14,c=WHT,b=True,a=PP_ALIGN.CENTER)

# Two branches
br_y=pool_y+Inches(0.85)
# Left: classification
card(s,Inches(0.6),br_y,Inches(5.8),Inches(2.7),border=GRN)
R(s,Inches(0.6),br_y,Inches(5.8),Inches(0.55),fill=GRN)
T(s,Inches(0.75),br_y+Inches(0.08),Inches(5.5),Inches(0.4),"分类分支 — Classification Branch",sz=15,c=WHT,b=True,a=PP_ALIGN.CENTER)
cls_br=[
("输入: 128维特征向量",13,G800,False),
("  ->  Linear(128->64) + ReLU + Dropout(0.25)",13,SEC,True),
("  ->  Linear(64->2)  输出层",13,PRI,True),
("  ->  输出: 2类 logits (特征不存在 vs 特征存在)",14,GRN,True),
("",8,TXT,False),
("4个特征的二分类标签:",13,TXT,True),
("  boundary:      [logit_smooth, logit_not_smooth]",12,G600,False),
("  calcification: [logit_no_cal, logit_cal]",12,G600,False),
("  shape:         [logit_regular, logit_irregular]",12,G600,False),
("  direction:     [logit_parallel, logit_not_parallel]",12,G600,False),
]
MT(s,Inches(0.8),br_y+Inches(0.65),Inches(5.4),Inches(1.9),cls_br)

# Right: regression
card(s,Inches(6.9),br_y,Inches(5.8),Inches(2.7),border=ACC)
R(s,Inches(6.9),br_y,Inches(5.8),Inches(0.55),fill=ACC)
T(s,Inches(7.05),br_y+Inches(0.08),Inches(5.5),Inches(0.4),"回归分支 — Regression Branch",sz=15,c=WHT,b=True,a=PP_ALIGN.CENTER)
reg_br=[
("输入: 128维特征向量",13,G800,False),
("  ->  Linear(128->64) + ReLU + Dropout(0.25)",13,SEC,True),
("  ->  Linear(64->4)  输出层",13,PRI,True),
("  ->  Sigmoid 激活 -> [xc, yc, w, h] 归一化坐标",14,ACC,True),
("",8,TXT,False),
("坐标说明:",13,TXT,True),
("  xc, yc: 边界框中心点 (归一化, 范围 [0,1])",12,G600,False),
("  w, h:   边界框宽高   (归一化, 范围 [0,1])",12,G600,False),
("",8,TXT,False),
("推理时:  bbox * [W, H, W, H] -> 像素坐标",12,HIL,True),
]
MT(s,Inches(7.1),br_y+Inches(0.65),Inches(5.4),Inches(1.9),reg_br)

PN(s,7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Loss Functions
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"损失函数设计","Multi-Task Loss — FocalLoss + DetectionLoss + Uncertainty Weighting")

# Master formula
R(s,Inches(0.5),Inches(1.4),Inches(12.3),Inches(0.7),fill=PRI)
T(s,Inches(0.7),Inches(1.42),Inches(11.9),Inches(0.66),
   "L_total = 1/(2*e^log_var_cls) * FocalLoss_BI-RADS  +  1/(2*e^log_var_det) * SUM w_f * [FocalLoss(cls|gamma_f) + 2.0*SmoothL1(bbox)]  +  log_var_cls + log_var_det\n"
   "其中 log_var_cls, log_var_det 为可学习的不确定性参数, w_f 为特征任务权重, gamma_f 为特征级Focal Gamma",
   sz=13,c=WHT,b=False,a=PP_ALIGN.CENTER)

# Three columns
cols=[
("FocalLoss 分类损失",GRN,[
("核心公式: FL(p_t) = -alpha * (1-p_t)^gamma * log(p_t)",13,TXT,True,PP_ALIGN.LEFT,Pt(5)),
("",6,TXT,False),
("BI-RADS分类配置:",14,PRI,True),
("  gamma=2.0, alpha=0.25, label_smoothing=0.1",13,G800,False),
("  类别权重: class_2=0.5, 其余=1.0",13,G800,False),
("",6,TXT,False),
("特征级 Focal Gamma (抑制假阳性):",14,PRI,True),
("  boundary:      gamma=2.5",13,G800,False),
("  calcification: gamma=2.0 (最低, 正负样本均衡)",13,G800,False),
("  shape:         gamma=3.0 (高, 正样本仅16.3%)",13,G800,False),
("  direction:     gamma=3.0 (高, 正样本仅16.6%)",13,G800,False),
("",6,TXT,False),
("gamma越高->对易分类样本损失压缩越强->迫使模型关注困难样本",13,HIL,True),
]),
("DetectionLoss 检测损失",ACC,[
("每个特征的检测损失分解:",14,PRI,True,PP_ALIGN.LEFT,Pt(5)),
("  L_det(f) = FocalLoss(cls_logits, gt_cls)",13,G800,False),
("           + 2.0 * SmoothL1(bbox_pred, gt_bbox)",13,G800,False),
("",6,TXT,False),
("分类部分:",14,PRI,True),
("  FocalLoss + 类别权重(0 vs 1)",13,G800,False),
("  权重策略: sqrt(1/freq) 温和平衡",13,G800,False),
("",6,TXT,False),
("回归部分:",14,PRI,True),
("  SmoothL1Loss (对异常值鲁棒)",13,G800,False),
("  bbox权重 x2.0 (强调定位精度)",13,G800,False),
("  仅对正样本计算bbox损失",13,G800,False),
("",6,TXT,False),
("特征任务权重 w_f:",14,PRI,True),
("  boundary:1.5 | calcification:1.0",13,G800,False),
("  shape:2.0    | direction:2.5 (最难)",13,G800,False),
]),
("不确定性加权 (Kendall 2018)",HIL,[
("可学习参数:",14,PRI,True,PP_ALIGN.LEFT,Pt(5)),
("  log_var_cls, log_var_det",13,G800,False),
("",6,TXT,False),
("自适应权重:",14,PRI,True),
("  precision = 1/(2*exp(log_var)+eps)",13,G800,False),
("  L_w = precision*L_raw + log_var",13,G800,False),
("",6,TXT,False),
("直觉解释:",14,PRI,True),
("  任务越困难->log_var变大",13,G800,False),
("  ->precision变小",13,G800,False),
("  ->该任务梯度贡献自动减小",13,G800,False),
("  ->避免困难任务主导训练方向",13,G800,False),
("",6,TXT,False),
("核心优势:",14,GRN,True),
("  -无需手动调节cls/det权重比",13,G800,False),
("  -训练中自适应探索最佳平衡点",13,G800,False),
("  -数学等价于高斯似然最大化",13,G800,False),
]),
]
for i,(title,clr,items) in enumerate(cols):
    cx=Inches(0.5+i*4.15)
    card(s,cx,Inches(2.4),Inches(3.95),Inches(4.8))
    R(s,cx,Inches(2.4),Inches(3.95),Inches(0.5),fill=clr)
    T(s,cx+Inches(0.1),Inches(2.45),Inches(3.75),Inches(0.4),title,sz=15,c=WHT,b=True)
    MT(s,cx+Inches(0.1),Inches(3.0),Inches(3.75),Inches(4.1),items)

PN(s,8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation Metrics
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"评价指标体系与最佳模型结果","Evaluation Metrics & Best Model Performance")

ST(s,Inches(0.6),Inches(1.4),"BI-RADS 分类评价指标")
brm=[
["指标","公式","最佳值"],
["Accuracy","correct / total",f"{best_epoch['val']['birads_acc']*100:.1f}%"if HAS_DATA else"--"],
["Per-class Precision","TP_c/(TP_c+FP_c)","各类别分别计算"],
["Per-class Recall","TP_c/(TP_c+FN_c)","各类别分别计算"],
["Per-class F1 Score","2*P_c*R_c/(P_c+R_c)","各类别分别计算"],
]
tbl(s,brm,Inches(0.6),Inches(2.0),Inches(5.5),Inches(2.0),[Inches(2.0),Inches(1.8),Inches(1.7)])

ST(s,Inches(6.8),Inches(1.4),"特征检测评价指标 (4个特征各自评估)")
featm=[
["指标","公式","钙化(最优)"if HAS_DATA else""],
["Accuracy","(TP+TN)/Total",f"{best_epoch['val']['calcification_cls_acc']*100:.1f}%"if HAS_DATA else"--"],
["Precision","TP/(TP+FP)",f"{best_epoch['val']['calcification_precision']*100:.1f}%"if HAS_DATA else"--"],
["Recall","TP/(TP+FN)",f"{best_epoch['val']['calcification_recall']*100:.1f}%"if HAS_DATA else"--"],
["F1 Score(主指标)","2PR/(P+R)",f"{best_epoch['val']['calcification_f1']*100:.1f}%"if HAS_DATA else"--"],
["Specificity","TN/(TN+FP)",f"{best_epoch['val']['calcification_specificity']*100:.1f}%"if HAS_DATA else"--"],
["IoU","Intersect/Union",f"{best_epoch['val']['calcification_iou']*100:.1f}%"if HAS_DATA else"--"],
["AP@0.5","Correct+IoU>0.5",f"{best_epoch['val']['calcification_det_ap@0.5']*100:.1f}%"if HAS_DATA else"--"],
]
tbl(s,featm,Inches(6.8),Inches(2.0),Inches(5.8),Inches(3.0),[Inches(1.8),Inches(1.7),Inches(2.3)])

ST(s,Inches(0.6),Inches(4.3),"模型选择策略")
sel=[
("* 最佳模型选择依据: 4个特征的 平均F1分数 (Mean F1) — 而非 BI-RADS 准确率",14,HIL,True,PP_ALIGN.LEFT,Pt(3)),
("  原因: 特征检测(尤其方向/形状)的假阳性问题远比BI-RADS分类更难优化，是模型性能的真正瓶颈",13,G800,False,PP_ALIGN.LEFT,Pt(6)),
("* 早停策略: patience=25轮 — 连续25轮Mean F1不提升则自动停止训练 (实际训练48轮停止)",14,HIL,True),
]
MT(s,Inches(0.6),Inches(4.8),Inches(12.3),Inches(1.0),sel)

if HAS_DATA and best_epoch:
    b=best_epoch["val"]
    ST(s,Inches(0.6),Inches(5.5),f"最佳模型详细结果 (Epoch {best_epoch['epoch']}, Mean F1 = {best_f1*100:.1f}%)")
    summary=[
    ("BI-RADS Acc",f"{b['birads_acc']*100:.1f}%",PRI),("Mean F1",f"{best_f1*100:.1f}%",ACC),
    ("钙化 F1",f"{b['calcification_f1']*100:.1f}%",GRN),("边界 F1",f"{b['boundary_f1']*100:.1f}%",SEC),
    ("形状 F1",f"{b['shape_f1']*100:.1f}%",RGBColor(0xF3,0x9C,0x12)),("方向 F1",f"{b['direction_f1']*100:.1f}%",HIL),
    ("钙化 IoU",f"{b['calcification_iou']*100:.1f}%",ACC),("钙化 AP@0.5",f"{b['calcification_det_ap@0.5']*100:.1f}%",GRN),
    ]
    for i,(label,val,clr) in enumerate(summary):
        sx=Inches(0.6+i*1.55)
        R(s,sx,Inches(6.1),Inches(1.4),Inches(0.7),fill=WHT,border=clr,bw=Pt(2))
        T(s,sx+Inches(0.05),Inches(6.12),Inches(1.3),Inches(0.25),label,sz=10,c=G600,a=PP_ALIGN.CENTER)
        T(s,sx+Inches(0.05),Inches(6.4),Inches(1.3),Inches(0.3),val,sz=18,c=clr,b=True,a=PP_ALIGN.CENTER)

PN(s,9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Training Strategy
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"训练策略与超参数完整配置","Training Strategy & Complete Hyperparameter Configuration")

ST(s,Inches(0.6),Inches(1.4),"完整超参数配置表")
hp=[
["超参数","配置值","设计思路"],
["优化器","AdamW","解耦权重衰减的Adam, 更好泛化"],
["基础学习率","3e-4","适中初始值, 配合warmup平滑启动"],
["骨干学习率","3e-5 (x0.1)","保护ImageNet预训练权重, 仅微调"],
["权重衰减","3e-3","适度L2正则化, 防止过拟合"],
["Batch Size","12","RTX 4090 24GB显存最大化利用"],
["最大Epochs","100","配合早停(patience=25)可提前结束"],
["输入尺寸","640x320","保持超声图像W:H=2:1宽高比"],
["LR调度器","CosineAnnealing","平滑衰减, 收敛路径优于阶梯衰减"],
["Warmup轮数","5 epochs","LR从0线性增长至目标值, 避免初期震荡"],
["标签平滑","0.1","缓解过拟合, 提升模型校准能力"],
["梯度裁剪","max_norm=3.0","防止多任务梯度冲突和梯度爆炸"],
["混合精度","AMP(FP16)","FP16前向/反向, FP32权重更新, 节省~40%显存"],
["MixUp概率","50%步骤","强正则化手段, 显著提升泛化能力"],
]
tbl(s,hp,Inches(0.6),Inches(1.9),Inches(6.5),Inches(5.2),[Inches(1.4),Inches(1.4),Inches(3.7)])

ST(s,Inches(7.5),Inches(1.4),"六大核心训练策略")
strategies=[
("1 分层学习率","骨干0.1x基础LR微调预训练特征，检测头和分类头全量LR从头训练。充分利用预训练知识，同时允许任务特定层快速适应。",ACC),
("2 学习率预热","前5个epoch LR从0线性增长到目标值，避免训练初期模型对数据分布未适应时的大梯度更新破坏预训练权重。",SEC),
("3 余弦退火","LR沿余弦曲线平滑衰减至接近0，相比阶梯衰减收敛路径更平滑，更易找到平坦最优解，泛化性能更好。",SEC),
("4 自动混合精度","前向/反向使用FP16加速，权重更新保持FP32全精度。显存减少~40%，训练速度提升~1.5x，精度无显著影响。",GRN),
("5 梯度裁剪","梯度L2范数超过3.0时等比例缩放。多任务训练中不同任务梯度量级差异大，裁剪防止某任务梯度主导更新方向。",HIL),
("6 MixUp增强","50%步骤执行: x=lambda*A+(1-lambda)*B，标签按lambda加权。作为强正则化有效缓解医学数据量有限时的过拟合问题。",ACC),
]
for i,(title,desc,clr) in enumerate(strategies):
    sy=Inches(1.8+i*0.9)
    card(s,Inches(7.5),sy,Inches(5.4),Inches(0.82))
    R(s,Inches(7.5),sy,Inches(0.06),Inches(0.82),fill=clr)
    T(s,Inches(7.7),sy+Inches(0.05),Inches(5.0),Inches(0.28),title,sz=13,c=clr,b=True)
    T(s,Inches(7.7),sy+Inches(0.38),Inches(5.0),Inches(0.4),desc,sz=10,c=G600)

PN(s,10)

# ═══════════════════════════════════════════════════════════════
# SLIDES 11-14 — Training Visualizations (2 images per slide + analysis)
# ═══════════════════════════════════════════════════════════════

# Helper for visualization slides — spacious layout: image left + analysis right
def vis_slide(slide_num, title, subtitle, images_spec):
    """
    images_spec: list of (img_filename, card_title, analysis_lines)
    Each slide has 2 rows, each row = [image (left 55%) + analysis text (right 40%)]
    """
    s=BS(); BG(s)
    TB(s,title,subtitle)

    for idx,(fn,ctitle,analysis) in enumerate(images_spec):
        row_y = Inches(1.4 + idx*2.95)
        # Left: image card
        img_w = Inches(5.8); img_h = Inches(2.65)
        card(s, Inches(0.4), row_y, img_w, img_h)
        T(s, Inches(0.55), row_y+Inches(0.06), img_w-Inches(0.3), Inches(0.28),
           ctitle, sz=13, c=PRI, b=True)
        img_path = IMG_DIR / fn
        if img_path.exists():
            try:
                s.shapes.add_picture(str(img_path), Inches(0.5), row_y+Inches(0.38),
                                     img_w-Inches(0.2), img_h-Inches(0.48))
            except: pass

        # Right: analysis panel
        an_w = Inches(6.5); an_h = Inches(2.65)
        card(s, Inches(6.5), row_y, an_w, an_h, border=ACC)
        R(s, Inches(6.5), row_y, Inches(0.06), an_h, fill=ACC)
        MT(s, Inches(6.8), row_y+Inches(0.2), an_w-Inches(0.4), an_h-Inches(0.4), analysis)
    PN(s, slide_num)

# --- Slide 11: Loss curves + Accuracy ---
vis_slide(11,"训练结果可视化 (1/5)","Training Visualization — Loss & Accuracy",[
("loss_curves.png", "损失曲线总览 (Loss Curves)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Total Loss: 训练与验证Loss同步持续下降，两者差距小，无明显过拟合",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("BIRADS Loss: 训练Loss下降平稳，验证Loss在Epoch 20后趋于稳定",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Detection Loss: 从~0.4降至~0.19，下降幅度最大；验证Loss略高于训练Loss",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("结论: 三条Loss曲线均健康收敛，未出现过拟合趋势，早停于Epoch 48合理",12,HIL,True),
]),
("accuracy.png", "BI-RADS 分类准确率曲线",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("训练Acc@1: 稳步上升至74.1%，表明模型在训练集上学习充分",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("验证Acc: 最终达到64.1%，绿色虚线标注最佳Epoch峰值",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Train-Val差距: 约10个百分点，存在一定程度的过拟合",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("建议: 可增加MixUp概率、增大Label Smoothing或降低模型容量来缩小gap",12,HIL,True),
]),
])

# --- Slide 12: IoU + Feature Metrics ---
vis_slide(12,"训练结果可视化 (2/5)","Training Visualization — IoU & Feature Classification",[
("iou_curves.png", "特征检测 IoU 曲线 (正样本BBox交并比)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("范围: 4个特征的IoU均收敛至0.46-0.49区间，定位精度整体一致",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("钙化(橙线)IoU最高约0.48，方向(红线)IoU略低约0.46",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("趋势: 训练初期快速提升(前10 epoch)，之后缓慢改善趋于平稳",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("提升空间: IoU距离理想值(>0.7)仍有差距，可尝试更强的检测头结构",12,HIL,True),
]),
("feature_cls_metrics.png", "特征分类完整指标 (P / R / F1 / Specificity)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Precision: 钙化(~0.69)远高于方向(~0.20)，方向假阳性问题是主要瓶颈",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Recall: 方向最高(~0.91)说明模型对方向正样本非常敏感，边界和钙化在0.73-0.79",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("F1 Score: 钙化(0.72)最优, 方向(0.33)最差 — 方向的高Recall被低Precision拖累",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("结论: 方向呈现典型的高Recall+低Precision模式=严重假阳性, 是性能最大瓶颈",12,HIL,True),
]),
])

# --- Slide 13: Detection loss breakdown + Feature cls acc ---
vis_slide(13,"训练结果可视化 (3/5)","Training Visualization — Detection Loss & Feature Accuracy",[
("det_breakdown.png", "特征检测 Loss 逐特征分解 (Cls Loss + BBox Loss)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("每个特征独立展示分类Loss(蓝线)和边界框Loss(红线)的逐epoch变化",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("boundary/shape的cls loss(蓝线)收敛值较高，反映正负样本不平衡导致的分类困难",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("direction的cls loss最低(~0.005)，说明分类过于偏向负样本(假阳性恰恰源于此)",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("所有bbox loss(红线)稳步下降至0.01-0.02，定位能力持续改善且4特征表现一致",12,G800,False),
]),
("feature_cls_acc.png", "特征二分类准确率对比 (4特征Accuracy)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("direction(红线)准确率最高约95%，但这是因为数据集中方向负样本占绝大多数",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("calcification(橙色)准确率约78%，因正负样本较均衡(~1:1)，准确率更具参考价值",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("boundary(蓝色)和shape(绿色)居中，约80%",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("重要提醒: 在不平衡数据上，Accuracy有欺骗性！必须结合P/R/F1综合评估",12,HIL,True),
]),
])

# --- Slide 14: Detection AP + LR Curve ---
vis_slide(14,"训练结果可视化 (4/5)","Training Visualization — Detection AP & Learning Rate",[
("detection_ap.png", "检测 AP@0.5 曲线 (类别正确 + IoU>0.5 的检出率)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("AP@0.5 综合衡量分类精度(类别正确)与定位精度(IoU>0.5)，比单独看更全面",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("钙化(橙线)AP@0.5最高约0.48，是综合表现最优的特征",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("方向(红线)AP@0.5仅约0.36，边界(蓝线)约0.41，形状(绿线)约0.40",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("改进: 针对方向/形状提高Focal Gamma至3.5+，引入OHEM困难样本挖掘",12,HIL,True),
]),
("lr_curve.png", "学习率调度曲线 (Log-scale Y轴)",[
("分析:",14,PRI,True,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("Log尺度Y轴清晰展示 CosineAnnealing + 5 epoch Linear Warmup 的完整过程",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("预热阶段: 前5个epoch，LR从0线性增长至3e-4，避免初期梯度震荡",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("退火阶段: Epoch 5-100沿余弦曲线平滑衰减，确保后期以极小LR(~1e-5)精细调优",12,G800,False,PP_ALIGN.LEFT,Pt(4)),
("",6,TXT,False),
("实际终止: 第48轮因早停触发，此时LR约1.77e-4，模型已收敛至最佳点附近",12,HIL,True),
]),
])

# --- Slide 15: Summary Dashboard + Confusion ---
s=BS(); BG(s)
TB(s,"训练结果可视化 (5/5) — 综合仪表盘","Training Summary Dashboard & Confusion Matrix")

# Summary image — left side
sp=IMG_DIR/"summary.png"
if sp.exists():
    try:
        s.shapes.add_picture(str(sp), Inches(0.2), Inches(1.5), Inches(8.3), Inches(5.6))
    except: pass

# Right: confusion matrix + analysis
cm_path=IMG_DIR/"confusion_matrices.png"
if cm_path.exists():
    try:
        s.shapes.add_picture(str(cm_path), Inches(8.8), Inches(1.5), Inches(4.0), Inches(3.0))
    except: pass
    T(s, Inches(8.8), Inches(4.6), Inches(4.0), Inches(0.25),
       "混淆矩阵 — 正负样本分类分布", sz=12, c=PRI, b=True, a=PP_ALIGN.CENTER)

# Key findings panel
ST(s, Inches(8.8), Inches(5.0), "关键发现")
if HAS_DATA and best_epoch:
    b=best_epoch["val"]
    findings=[
    (f"BI-RADS Acc: {b['birads_acc']*100:.1f}%  |  Mean F1: {best_f1*100:.1f}%",13,PRI,True,PP_ALIGN.LEFT,Pt(4)),
    (f"钙化最优: F1={b['calcification_f1']*100:.1f}%, P={b['calcification_precision']*100:.1f}%, R={b['calcification_recall']*100:.1f}%",12,GRN,True,PP_ALIGN.LEFT,Pt(4)),
    (f"方向最难: F1={b['direction_f1']*100:.1f}%, P={b['direction_precision']*100:.1f}%, R={b['direction_recall']*100:.1f}%",12,HIL,True,PP_ALIGN.LEFT,Pt(4)),
    (f"形状次难: F1={b['shape_f1']*100:.1f}%, P={b['shape_precision']*100:.1f}%, R={b['shape_recall']*100:.1f}%",12,RGBColor(0xF3,0x9C,0x12),True,PP_ALIGN.LEFT,Pt(4)),
    ("",6,TXT,False),
    ("改进: 提高方向/形状Focal Gamma | OHEM | 更强检测头 | 增加样本",12,G800,False),
    ]
    MT(s, Inches(8.8), Inches(5.4), Inches(4.2), Inches(2.0), findings)

PN(s,15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Inference
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"推理流程与输出格式","Inference Pipeline & Output Specification — infer.py")

# Flow
flow_row(s,Inches(0.3),Inches(1.45),[
("1.加载测试集\nTestDataset","BIRADS分类测试集\n+特征检测测试集\n保留原始图像尺寸\n用于bbox反归一化",PRI),
("2.图像预处理\n归一化","Resize->640x320\nImageNet标准化\nmean=[0.485,0.456,0.406]\nstd=[0.229,0.224,0.225]",SEC),
("3.加载模型\nload_checkpoint","从best.pt恢复\n模型权重+配置\n@torch.no_grad()\n无梯度推理模式",SEC),
("4.模型前向\nmodel(images)","共享骨干+FPN\n5个头并行输出\nbirads_logits+\nfeature_outputs",ACC),
("5.后处理\n解析结果","argmax->类别索引\nSigmoid bbox->[0,1]\n反归一化->像素坐标\n映射类别名称",HIL),
],bw=Inches(2.3),bh=Inches(1.85))

# Two output formats
ST(s,Inches(0.6),Inches(3.65),"双路输出 JSON 格式")
# Left: Classification
card(s,Inches(0.6),Inches(4.2),Inches(5.8),Inches(2.9),border=PRI)
R(s,Inches(0.6),Inches(4.2),Inches(5.8),Inches(0.55),fill=PRI)
T(s,Inches(0.75),Inches(4.25),Inches(5.5),Inches(0.4),"class_result.json — BI-RADS 分类输出",sz=15,c=WHT,b=True,a=PP_ALIGN.CENTER)
cls_json=[
("[",12,G600,False),
('  { "image_id": "P00001",',12,G800,False),
('    "ground_truth": "4A类",',12,G800,False),
('    "predicted_birads": "4A类"',12,GRN,True),
("  },",12,G600,False),
('  { "image_id": "P00002",',12,G800,False),
('    "ground_truth": "3类",',12,G800,False),
('    "predicted_birads": "4A类"',12,HIL,True),
("  }",12,G600,False),
("]",12,G600,False),
]
MT(s,Inches(0.8),Inches(4.75),Inches(5.4),Inches(2.2),cls_json)

# Right: Feature detection
card(s,Inches(6.9),Inches(4.2),Inches(5.8),Inches(2.9),border=HIL)
R(s,Inches(6.9),Inches(4.2),Inches(5.8),Inches(0.55),fill=HIL)
T(s,Inches(7.05),Inches(4.25),Inches(5.5),Inches(0.4),"future_result.json — 特征检测输出",sz=15,c=WHT,b=True,a=PP_ALIGN.CENTER)
feat_json=[
("[",12,G600,False),
('  { "image_id": "P00001",',12,G800,False),
('    "boundary": {',12,G800,False),
('      "class": "smooth",',12,ACC,True),
('      "bbox": [120.5, 85.3, 45.2, 38.7] },',12,ACC,True),
('    "calcification": {',12,G800,False),
('      "class": "calcification",',12,GRN,True),
('      "bbox": [200.1, 150.8, 30.5, 25.2] },',12,GRN,True),
('    "shape": { "class": "regular", "bbox": [...] },',12,G800,False),
('    "direction": { "class": "parallel", "bbox": [...] }',12,G800,False),
("  }",12,G600,False),
("]",12,G600,False),
("注: bbox坐标为反归一化后的像素坐标 [xc, yc, w, h]",11,G600,False),
]
MT(s,Inches(7.1),Inches(4.75),Inches(5.4),Inches(2.2),feat_json)

PN(s,16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Confusion Matrix & Summary
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s)
TB(s,"训练结果 — 混淆矩阵与总结分析","Confusion Matrices & Summary Analysis")

# Confusion matrix if exists
cm_path=IMG_DIR/"confusion_matrices.png"
if cm_path.exists():
    try:
        s.shapes.add_picture(str(cm_path),Inches(0.4),Inches(1.5),Inches(6.5),Inches(4.5))
    except: pass
    T(s,Inches(0.5),Inches(6.1),Inches(6.3),Inches(0.3),"混淆矩阵 — 各特征正负样本分类详细分布",sz=13,c=PRI,b=True)

# Right side: overall summary
ST(s,Inches(7.3),Inches(1.5),"模型性能总结 (Epoch 29)")
if HAS_DATA and best_epoch:
    b=best_epoch["val"]
    perf_table=[
    ["特征","类别准确率","Precision","Recall","F1 Score","IoU","AP@0.5"],
    ["boundary",f"{b['boundary_cls_acc']*100:.1f}%",f"{b['boundary_precision']*100:.1f}%",f"{b['boundary_recall']*100:.1f}%",f"{b['boundary_f1']*100:.1f}%",f"{b['boundary_iou']*100:.1f}%",f"{b['boundary_det_ap@0.5']*100:.1f}%"],
    ["calcification",f"{b['calcification_cls_acc']*100:.1f}%",f"{b['calcification_precision']*100:.1f}%",f"{b['calcification_recall']*100:.1f}%",f"{b['calcification_f1']*100:.1f}%",f"{b['calcification_iou']*100:.1f}%",f"{b['calcification_det_ap@0.5']*100:.1f}%"],
    ["shape",f"{b['shape_cls_acc']*100:.1f}%",f"{b['shape_precision']*100:.1f}%",f"{b['shape_recall']*100:.1f}%",f"{b['shape_f1']*100:.1f}%",f"{b['shape_iou']*100:.1f}%",f"{b['shape_det_ap@0.5']*100:.1f}%"],
    ["direction",f"{b['direction_cls_acc']*100:.1f}%",f"{b['direction_precision']*100:.1f}%",f"{b['direction_recall']*100:.1f}%",f"{b['direction_f1']*100:.1f}%",f"{b['direction_iou']*100:.1f}%",f"{b['direction_det_ap@0.5']*100:.1f}%"],
    ]
    tbl(s,perf_table,Inches(7.3),Inches(2.1),Inches(5.6),Inches(2.2),
        [Inches(1.1),Inches(0.85),Inches(0.85),Inches(0.75),Inches(0.7),Inches(0.65),Inches(0.7)])
    # Overall
    T(s,Inches(7.3),Inches(4.5),Inches(5.6),Inches(0.1),f"BI-RADS 分类准确率: {b['birads_acc']*100:.1f}%  |  平均F1: {best_f1*100:.1f}%  |  共训练48轮, 最佳epoch=29",sz=12,c=PRI,b=True)
# Key insight
ST(s,Inches(7.3),Inches(4.8),"核心发现")
insights=[
("1. 钙化检测最优: 正负样本比例均衡(~1:1), F1达71.6%",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("2. 方向检测最难: 高Recall(90.9%)但低Precision(19.6%), 严重假阳性",13,HIL,True,PP_ALIGN.LEFT,Pt(3)),
("3. 形状类似方向: Recall 74.0% vs Precision 35.2%, 也偏向过预测",13,HIL,True,PP_ALIGN.LEFT,Pt(3)),
("4. 边界处于中间: F1 55.0%, 平衡性相对较好",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("5. BI-RADS分类 64.1%: 6类分级有一定难度, 尤其相邻等级易混淆",13,G800,False,PP_ALIGN.LEFT,Pt(3)),
("6. 定位精度一致: 4个特征IoU均在46-49%, 检测头结构还有提升空间",13,G800,False),
]
MT(s,Inches(7.3),Inches(5.25),Inches(5.6),Inches(2.0),insights)

PN(s,17)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — Summary & Future
# ═══════════════════════════════════════════════════════════════
s=BS(); BG(s,PRI)
R(s,Inches(0),Inches(0),SW,Inches(0.06),fill=HIL)
T(s,Inches(0.8),Inches(0.5),Inches(11),Inches(0.6),"总结与展望",sz=36,c=WHT,b=True)
R(s,Inches(0.8),Inches(1.1),Inches(1.2),Inches(0.04),fill=HIL)

T(s,Inches(0.8),Inches(1.5),Inches(5.5),Inches(0.4),"核心贡献",sz=22,c=HIL,b=True)
cl=[
("* 多任务统一框架",18,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  单模型同时完成BI-RADS 6类分级+4项特征检测+区域定位",13.5,G400,False,PP_ALIGN.LEFT,Pt(10)),
("* 针对性损失函数设计",18,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  不确定性加权+特征级Focal Gamma+sqrt(1/freq)类别权重",13.5,G400,False,PP_ALIGN.LEFT,Pt(10)),
("* 完善的工程化实现",18,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  数据校验->训练->评估->可视化->推理 全流程覆盖",13.5,G400,False,PP_ALIGN.LEFT,Pt(10)),
("* 丰富的评价体系",18,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  Acc/P/R/F1/Specificity/IoU/AP@0.5 七维指标",13.5,G400,False),
]
MT(s,Inches(0.8),Inches(2.1),Inches(5.5),Inches(3.5),cl)

T(s,Inches(7.0),Inches(1.5),Inches(5.5),Inches(0.4),"未来方向",sz=22,c=HIL,b=True)
fl=[
("1.注意力机制增强",17,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  添加CBAM/SE/Transformer模块提升特征表达",13,G400,False,PP_ALIGN.LEFT,Pt(8)),
("2.更强的检测头设计",17,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  Anchor-based/DETR-style Transformer检测头",13,G400,False,PP_ALIGN.LEFT,Pt(8)),
("3.假阳性专项治理",17,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  针对方向/形状的OHEM+Focal Loss精细调优",13,G400,False,PP_ALIGN.LEFT,Pt(8)),
("4.模型轻量化与部署",17,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  知识蒸馏+INT8量化->边缘设备实时推理",13,G400,False,PP_ALIGN.LEFT,Pt(8)),
("5.多模态与外部验证",17,WHT,True,PP_ALIGN.LEFT,Pt(6)),
("  结合临床文本+多中心数据验证泛化能力",13,G400,False),
]
MT(s,Inches(7.0),Inches(2.1),Inches(5.5),Inches(3.5),fl)

R(s,Inches(0),SH-Inches(0.06),SW,Inches(0.06),fill=HIL)
T(s,Inches(0.5),SH-Inches(0.7),Inches(12),Inches(0.5),"谢谢！  Thank You!",sz=30,c=WHT,b=True,a=PP_ALIGN.CENTER)
PN(s,18)

# ── Save ──
prs.save(str(PPT_PATH))
print(f"DONE -> {PPT_PATH}")
print(f"Slides: {len(prs.slides)}, Data: {HAS_DATA} ({total_epochs} epochs)")
if HAS_DATA: print(f"Best epoch: {best_epoch['epoch']}, Mean F1: {best_f1:.4f}")
